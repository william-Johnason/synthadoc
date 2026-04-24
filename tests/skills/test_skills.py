# SPDX-License-Identifier: AGPL-3.0-or-later
# Copyright (C) 2026 Paul Chen / axoviq.com
import pytest
from pathlib import Path


def test_all_builtin_skills_registered():
    from synthadoc.agents.skill_agent import SkillAgent
    agent = SkillAgent()
    names = [s.name for s in agent.list_skills()]
    for expected in ("pdf", "url", "markdown", "docx", "xlsx", "image", "pptx"):
        assert expected in names


def test_detect_skill_by_extension():
    from synthadoc.agents.skill_agent import SkillAgent
    agent = SkillAgent()
    assert agent.detect_skill("paper.pdf").name == "pdf"
    assert agent.detect_skill("notes.md").name == "markdown"
    assert agent.detect_skill("report.docx").name == "docx"
    assert agent.detect_skill("data.xlsx").name == "xlsx"
    assert agent.detect_skill("photo.png").name == "image"
    assert agent.detect_skill("deck.pptx").name == "pptx"


def test_detect_skill_by_url_prefix():
    from synthadoc.agents.skill_agent import SkillAgent
    agent = SkillAgent()
    assert agent.detect_skill("https://example.com/page").name == "url"
    assert agent.detect_skill("http://example.com").name == "url"


def test_unknown_extension_raises():
    from synthadoc.agents.skill_agent import SkillAgent, SkillNotFoundError
    agent = SkillAgent()
    with pytest.raises(SkillNotFoundError):
        agent.detect_skill("file.xyz")


def test_tier1_metadata_always_available():
    from synthadoc.agents.skill_agent import SkillAgent
    agent = SkillAgent()
    for meta in agent.list_skills():
        assert meta.name
        assert meta.description


@pytest.mark.asyncio
async def test_markdown_skill_extracts(tmp_path):
    from synthadoc.skills.markdown.scripts.main import MarkdownSkill
    f = tmp_path / "note.md"
    f.write_text("# Hello\n\nWorld content here.", encoding="utf-8")
    result = await MarkdownSkill().extract(str(f))
    assert "Hello" in result.text and "World" in result.text


@pytest.mark.asyncio
async def test_url_skill_fetches(tmp_path):
    import respx, httpx
    from synthadoc.skills.url.scripts.main import UrlSkill
    with respx.mock:
        respx.get("https://example.com/").mock(
            return_value=httpx.Response(200,
                text="<html><body><p>Hello from web</p></body></html>")
        )
        result = await UrlSkill().extract("https://example.com/")
    assert "Hello from web" in result.text


@pytest.mark.asyncio
async def test_url_skill_raises_on_404():
    import respx, httpx
    from synthadoc.skills.url.scripts.main import UrlSkill
    with respx.mock:
        respx.get("https://example.com/missing").mock(return_value=httpx.Response(404))
        with pytest.raises(Exception, match="404"):
            await UrlSkill().extract("https://example.com/missing")


@pytest.mark.asyncio
async def test_url_skill_raises_on_connection_error():
    import respx, httpx
    from synthadoc.skills.url.scripts.main import UrlSkill
    with respx.mock:
        respx.get("https://unreachable.example/").mock(side_effect=httpx.ConnectError("refused"))
        with pytest.raises(httpx.ConnectError):
            await UrlSkill().extract("https://unreachable.example/")


@pytest.mark.asyncio
async def test_url_skill_raises_domain_blocked_on_403():
    """403/401/429 responses must raise DomainBlockedException, not HTTPStatusError."""
    import respx, httpx
    from synthadoc.skills.url.scripts.main import UrlSkill
    from synthadoc.errors import DomainBlockedException
    with respx.mock:
        respx.get("https://blocked.example/page").mock(return_value=httpx.Response(403))
        with pytest.raises(DomainBlockedException) as exc_info:
            await UrlSkill().extract("https://blocked.example/page")
    assert exc_info.value.domain == "blocked.example"
    assert exc_info.value.status_code == 403


@pytest.mark.asyncio
async def test_url_skill_extracts_pdf_via_content_type():
    """URLs served with application/pdf content-type route to PDF extraction."""
    import respx, httpx
    from unittest.mock import MagicMock, patch
    from synthadoc.skills.url.scripts.main import UrlSkill
    from synthadoc.skills.base import ExtractedContent

    pdf_bytes = b"%PDF-1.4 fake"
    with respx.mock:
        respx.get("https://example.com/doc.pdf").mock(
            return_value=httpx.Response(200, content=pdf_bytes,
                headers={"content-type": "application/pdf"})
        )
        fake_result = ExtractedContent(text="PDF text here", source_path="https://example.com/doc.pdf",
                                       metadata={"url": "https://example.com/doc.pdf", "pages": 1})
        with patch.object(UrlSkill, "_extract_pdf_response", return_value=fake_result):
            result = await UrlSkill().extract("https://example.com/doc.pdf")
    assert result.text == "PDF text here"


def test_url_skill_extract_pdf_response_pypdf_success(tmp_path):
    """_extract_pdf_response returns pypdf text when pypdf succeeds."""
    from unittest.mock import MagicMock, patch
    from synthadoc.skills.url.scripts.main import UrlSkill

    skill = UrlSkill()
    pdf_bytes = b"%PDF-1.4 content"

    mock_page = MagicMock()
    mock_page.extract_text.return_value = "Extracted PDF text from page."
    mock_reader = MagicMock()
    mock_reader.pages = [mock_page]

    with patch("pypdf.PdfReader", return_value=mock_reader):
        result = skill._extract_pdf_response(pdf_bytes, "https://example.com/test.pdf")

    assert "Extracted PDF text from page." in result.text
    assert result.metadata["pages"] == 1


def test_url_skill_extract_pdf_response_pypdf_empty_falls_back_to_pdfminer(tmp_path):
    """When pypdf yields no text, _extract_pdf_response falls back to pdfminer."""
    from unittest.mock import MagicMock, patch
    from synthadoc.skills.url.scripts.main import UrlSkill

    skill = UrlSkill()
    pdf_bytes = b"%PDF-1.4 empty"

    mock_page = MagicMock()
    mock_page.extract_text.return_value = ""  # pypdf gives nothing
    mock_reader = MagicMock()
    mock_reader.pages = [mock_page]

    with patch("pypdf.PdfReader", return_value=mock_reader), \
         patch("pdfminer.high_level.extract_text", return_value="pdfminer extracted text"):
        result = skill._extract_pdf_response(pdf_bytes, "https://example.com/test.pdf")

    assert "pdfminer extracted text" in result.text


def test_url_skill_extract_pdf_response_both_fail_returns_empty(tmp_path):
    """When both pypdf and pdfminer fail, _extract_pdf_response returns empty ExtractedContent."""
    from unittest.mock import patch
    from synthadoc.skills.url.scripts.main import UrlSkill

    skill = UrlSkill()
    pdf_bytes = b"%PDF-1.4 broken"

    with patch("pypdf.PdfReader", side_effect=RuntimeError("corrupt")), \
         patch("pdfminer.high_level.extract_text", side_effect=RuntimeError("also broken")):
        result = skill._extract_pdf_response(pdf_bytes, "https://example.com/bad.pdf")

    assert result.text == ""
    assert result.metadata.get("pages") == 0


def test_pip_skills_loaded_from_entry_points(tmp_wiki):
    """Skills registered via entry_points('synthadoc.skills') are auto-discovered."""
    import yaml
    from unittest.mock import patch
    from synthadoc.agents.skill_agent import SkillAgent

    skill_dir = tmp_wiki / "skills" / "custom_pip"
    (skill_dir / "scripts").mkdir(parents=True)
    fm = {
        "name": "custom_pip", "version": "1.0", "description": "pip skill",
        "entry": {"script": "scripts/main.py", "class": "CustomPipSkill"},
        "triggers": {"extensions": [".xyz"], "intents": []}, "requires": [],
    }
    (skill_dir / "SKILL.md").write_text(f"---\n{yaml.dump(fm)}---\n", encoding="utf-8")
    (skill_dir / "scripts" / "main.py").write_text(
        "from synthadoc.skills.base import BaseSkill, ExtractedContent\n"
        "class CustomPipSkill(BaseSkill):\n"
        "    async def extract(self, s): return ExtractedContent('', s, {})\n",
        encoding="utf-8",
    )
    with patch("synthadoc.agents.skill_agent._entry_point_skill_dirs",
               return_value=[skill_dir]):
        agent = SkillAgent(wiki_root=tmp_wiki)
    assert "custom_pip" in [s.name for s in agent.list_skills()]


def test_local_skill_loaded_from_skills_folder(tmp_path):
    """A skill folder dropped in wiki/skills/ is discovered and registered."""
    import yaml
    from synthadoc.agents.skill_agent import SkillAgent

    (tmp_path / "skills").mkdir()
    (tmp_path / ".synthadoc").mkdir()
    skill_dir = tmp_path / "skills" / "local_csv"
    (skill_dir / "scripts").mkdir(parents=True)
    fm = {
        "name": "local_csv", "version": "1.0", "description": "local",
        "entry": {"script": "scripts/main.py", "class": "LocalCsvSkill"},
        "triggers": {"extensions": [".tsv"], "intents": []}, "requires": [],
    }
    (skill_dir / "SKILL.md").write_text(f"---\n{yaml.dump(fm)}---\n", encoding="utf-8")
    (skill_dir / "scripts" / "main.py").write_text(
        "from synthadoc.skills.base import BaseSkill, ExtractedContent\n"
        "class LocalCsvSkill(BaseSkill):\n"
        "    async def extract(self, s): return ExtractedContent('', s, {})\n",
        encoding="utf-8",
    )
    agent = SkillAgent(wiki_root=tmp_path)
    assert "local_csv" in [s.name for s in agent.list_skills()]


def test_local_skill_overrides_builtin(tmp_path):
    """A local skill folder with the same name as a built-in takes precedence."""
    import yaml
    from synthadoc.agents.skill_agent import SkillAgent

    (tmp_path / "skills").mkdir()
    (tmp_path / ".synthadoc").mkdir()
    skill_dir = tmp_path / "skills" / "pdf"
    (skill_dir / "scripts").mkdir(parents=True)
    fm = {
        "name": "pdf", "version": "1.0", "description": "override",
        "entry": {"script": "scripts/main.py", "class": "OverridePdfSkill"},
        "triggers": {"extensions": [".pdf"], "intents": []}, "requires": [],
    }
    (skill_dir / "SKILL.md").write_text(f"---\n{yaml.dump(fm)}---\n", encoding="utf-8")
    (skill_dir / "scripts" / "main.py").write_text(
        "from synthadoc.skills.base import BaseSkill, ExtractedContent\n"
        "class OverridePdfSkill(BaseSkill):\n"
        "    async def extract(self, s): return ExtractedContent('custom', s, {})\n",
        encoding="utf-8",
    )
    agent = SkillAgent(wiki_root=tmp_path)
    assert agent.detect_skill("doc.pdf").description == "override"


def test_local_skill_with_syntax_error_is_skipped(tmp_path, caplog):
    """A skill folder with a broken SKILL.md is logged and skipped."""
    import logging
    from synthadoc.agents.skill_agent import SkillAgent

    (tmp_path / "skills").mkdir()
    (tmp_path / ".synthadoc").mkdir()
    bad_dir = tmp_path / "skills" / "broken"
    bad_dir.mkdir()
    (bad_dir / "SKILL.md").write_text("this is not valid yaml: {{{", encoding="utf-8")
    with caplog.at_level(logging.WARNING):
        agent = SkillAgent(wiki_root=tmp_path)
    assert any("broken" in r.message for r in caplog.records)
    assert "pdf" in [s.name for s in agent.list_skills()]   # built-ins still registered


def test_local_file_not_subclassing_baseskill_is_skipped(tmp_path, caplog):
    """A skill folder with no SKILL.md is silently ignored."""
    import logging
    from synthadoc.agents.skill_agent import SkillAgent

    (tmp_path / "skills").mkdir()
    (tmp_path / ".synthadoc").mkdir()
    bad_dir = tmp_path / "skills" / "notaskill"
    bad_dir.mkdir()
    # no SKILL.md — not a skill folder
    with caplog.at_level(logging.WARNING):
        agent = SkillAgent(wiki_root=tmp_path)
    assert "notaskill" not in [s.name for s in agent.list_skills()]


@pytest.mark.asyncio
async def test_pdf_skill_cjk_fallback(tmp_path, monkeypatch):
    """PdfSkill falls back to pdfminer when pypdf yields too little text."""
    from synthadoc.skills.pdf.scripts.main import PdfSkill
    import synthadoc.skills.pdf.scripts.main as pdf_mod

    # Simulate pypdf returning nearly nothing (typical CJK font failure)
    def fake_pypdf(self, source):
        return ("x", 5)  # 1 char for 5 pages — below threshold

    monkeypatch.setattr(PdfSkill, "_extract_pypdf", fake_pypdf)

    # Simulate pdfminer returning the real Chinese text
    def fake_pdfminer(self, source):
        return "人工智能是计算机科学的一个重要分支。"

    monkeypatch.setattr(PdfSkill, "_extract_pdfminer", fake_pdfminer)

    dummy_pdf = tmp_path / "test.pdf"
    dummy_pdf.write_bytes(b"%PDF-1.4")  # minimal placeholder
    result = await PdfSkill().extract(str(dummy_pdf))
    assert "人工智能" in result.text


@pytest.mark.asyncio
async def test_pdf_skill_no_fallback_when_pypdf_succeeds(tmp_path, monkeypatch):
    """PdfSkill does NOT call pdfminer when pypdf returns enough text."""
    from synthadoc.skills.pdf.scripts.main import PdfSkill

    def fake_pypdf(self, source):
        return ("A" * 500, 3)  # well above threshold

    monkeypatch.setattr(PdfSkill, "_extract_pypdf", fake_pypdf)

    called = []
    def fake_pdfminer(self, source):
        called.append(True)
        return "pdfminer text"

    monkeypatch.setattr(PdfSkill, "_extract_pdfminer", fake_pdfminer)

    dummy_pdf = tmp_path / "test.pdf"
    dummy_pdf.write_bytes(b"%PDF-1.4")
    await PdfSkill().extract(str(dummy_pdf))
    assert not called, "pdfminer should not be called when pypdf yields sufficient text"


def test_tier2_instantiates_only_on_demand():
    """get_skill() should not be called during SkillAgent init — only on demand."""
    from synthadoc.agents.skill_agent import SkillAgent
    agent = SkillAgent()
    metas = agent.list_skills()
    assert all(hasattr(m, "name") for m in metas)
    skill = agent.get_skill("pdf")
    assert skill is not None


def test_tier3_resource_loaded_lazily(tmp_path):
    """get_resource() loads a file from the skill's resources/ dir on first call only."""
    from synthadoc.skills.base import BaseSkill, ExtractedContent, SkillMeta

    resources_dir = tmp_path / "resources"
    resources_dir.mkdir()
    (resources_dir / "prompt.md").write_text("# Vision prompt", encoding="utf-8")

    class ResourceSkill(BaseSkill):
        meta = SkillMeta(name="res", description="d", extensions=[".res"])
        async def extract(self, source): return ExtractedContent("", source, {})

    skill = ResourceSkill()
    skill._resources_dir = resources_dir
    text = skill.get_resource("prompt.md")
    assert text == "# Vision prompt"
    # second call uses cache — file can be deleted
    (resources_dir / "prompt.md").unlink()
    assert skill.get_resource("prompt.md") == "# Vision prompt"


def test_triggers_dataclass():
    from synthadoc.skills.base import Triggers
    t = Triggers(extensions=[".pdf"], intents=["research paper"])
    assert ".pdf" in t.extensions
    assert "research paper" in t.intents


def test_skillmeta_with_triggers():
    from synthadoc.skills.base import SkillMeta, Triggers
    m = SkillMeta(
        name="pdf", version="1.0", description="PDF skill",
        entry_script="scripts/main.py", entry_class="PdfSkill",
        triggers=Triggers(extensions=[".pdf"], intents=["document"]),
        requires=["pypdf"],
    )
    assert m.triggers.extensions == [".pdf"]
    assert m.version == "1.0"


def test_skillmeta_backwards_compat_extensions():
    """Old-style SkillMeta(extensions=[...]) still works without triggers."""
    from synthadoc.skills.base import SkillMeta
    m = SkillMeta(name="old", description="old skill", extensions=[".old"])
    assert ".old" in m.triggers.extensions


def test_get_resource_searches_assets_then_references(tmp_path):
    from synthadoc.skills.base import BaseSkill, ExtractedContent, SkillMeta

    (tmp_path / "assets").mkdir()
    (tmp_path / "assets" / "config.json").write_text('{"key": "val"}', encoding="utf-8")
    (tmp_path / "references").mkdir()
    (tmp_path / "references" / "notes.md").write_text("# Notes", encoding="utf-8")

    class TestSkill(BaseSkill):
        meta = SkillMeta(name="t", description="t", extensions=[".t"])
        async def extract(self, source): return ExtractedContent("", source, {})

    skill = TestSkill()
    skill.skill_dir = tmp_path
    assert '{"key": "val"}' in skill.get_resource("config.json")
    assert "# Notes" in skill.get_resource("notes.md")


def test_get_resource_cache_avoids_second_read(tmp_path):
    from synthadoc.skills.base import BaseSkill, ExtractedContent, SkillMeta

    (tmp_path / "assets").mkdir()
    f = tmp_path / "assets" / "data.txt"
    f.write_text("original", encoding="utf-8")

    class TestSkill(BaseSkill):
        meta = SkillMeta(name="t2", description="t2", extensions=[".t2"])
        async def extract(self, source): return ExtractedContent("", source, {})

    skill = TestSkill()
    skill.skill_dir = tmp_path
    skill.get_resource("data.txt")
    f.write_text("modified", encoding="utf-8")
    assert skill.get_resource("data.txt") == "original"  # served from cache


# ── PPTX skill ────────────────────────────────────────────────────────────────

@pytest.mark.asyncio
async def test_pptx_skill_extracts_slide_text(tmp_path):
    """PptxSkill extracts title and body text from each slide."""
    from pptx import Presentation
    from pptx.util import Inches
    from synthadoc.skills.pptx.scripts.main import PptxSkill

    path = tmp_path / "deck.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]  # Title and Content

    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "Intro"
    slide.placeholders[1].text = "Hello world"

    prs.save(str(path))

    result = await PptxSkill().extract(str(path))
    assert "Intro" in result.text
    assert "Hello world" in result.text
    assert result.metadata["slides"] == 1


@pytest.mark.asyncio
async def test_pptx_skill_includes_speaker_notes(tmp_path):
    """PptxSkill appends speaker notes when present."""
    from pptx import Presentation
    from synthadoc.skills.pptx.scripts.main import PptxSkill

    path = tmp_path / "notes.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "Slide One"
    slide.notes_slide.notes_text_frame.text = "Remember to demo this"
    prs.save(str(path))

    result = await PptxSkill().extract(str(path))
    assert "Remember to demo this" in result.text


@pytest.mark.asyncio
async def test_pptx_skill_empty_presentation(tmp_path):
    """PptxSkill handles a presentation with no slides without error."""
    from pptx import Presentation
    from synthadoc.skills.pptx.scripts.main import PptxSkill

    path = tmp_path / "empty.pptx"
    Presentation().save(str(path))

    result = await PptxSkill().extract(str(path))
    assert result.text == ""
    assert result.metadata["slides"] == 0


@pytest.mark.asyncio
async def test_pdf_skill_invalid_file_raises_descriptive_error(tmp_path):
    """PdfSkill raises a clear ValueError for corrupt or non-PDF files."""
    from synthadoc.skills.pdf.scripts.main import PdfSkill
    bad = tmp_path / "bad.pdf"
    bad.write_bytes(b"this is not a pdf")
    with pytest.raises(ValueError, match="Cannot read.*PDF"):
        await PdfSkill().extract(str(bad))


@pytest.mark.asyncio
async def test_docx_skill_invalid_file_raises_descriptive_error(tmp_path):
    """DocxSkill raises a clear ValueError for corrupt or non-DOCX files."""
    from synthadoc.skills.docx.scripts.main import DocxSkill
    bad = tmp_path / "bad.docx"
    bad.write_bytes(b"this is not a docx")
    with pytest.raises(ValueError, match="Cannot read.*Word"):
        await DocxSkill().extract(str(bad))


@pytest.mark.asyncio
async def test_xlsx_skill_invalid_file_raises_descriptive_error(tmp_path):
    """XlsxSkill raises a clear ValueError for corrupt or non-XLSX files."""
    from synthadoc.skills.xlsx.scripts.main import XlsxSkill
    bad = tmp_path / "bad.xlsx"
    bad.write_bytes(b"this is not an xlsx")
    with pytest.raises(ValueError, match="Cannot read.*Excel"):
        await XlsxSkill().extract(str(bad))


@pytest.mark.asyncio
async def test_pptx_skill_invalid_file_raises_descriptive_error(tmp_path):
    """PptxSkill raises a clear ValueError for corrupt or non-PPTX files."""
    from synthadoc.skills.pptx.scripts.main import PptxSkill

    bad = tmp_path / "bad.pptx"
    bad.write_bytes(b"this is not a zip file")

    with pytest.raises(ValueError, match="Cannot read.*PowerPoint"):
        await PptxSkill().extract(str(bad))


def test_pptx_detected_by_intent():
    """Intent keywords route to pptx skill."""
    from synthadoc.agents.skill_agent import SkillAgent
    agent = SkillAgent()
    assert agent.detect_skill("powerpoint").name == "pptx"


@pytest.mark.asyncio
async def test_xlsx_skill_extracts_csv(tmp_path):
    """XlsxSkill reads .csv files via the csv path (not openpyxl)."""
    from synthadoc.skills.xlsx.scripts.main import XlsxSkill
    csv_file = tmp_path / "data.csv"
    csv_file.write_text("name,age\nAlice,30\nBob,25\n", encoding="utf-8")
    result = await XlsxSkill().extract(str(csv_file))
    assert "Alice" in result.text
    assert "Bob" in result.text
    assert "name" in result.text


@pytest.mark.asyncio
async def test_pdf_skill_pdfminer_fallback_used_for_low_yield(tmp_path):
    """PdfSkill falls back to pdfminer when pypdf yields too little text."""
    from unittest.mock import patch, MagicMock
    from synthadoc.skills.pdf.scripts.main import PdfSkill

    skill = PdfSkill()
    # Simulate pypdf returning sparse text (below threshold) for a 2-page doc
    with patch.object(skill, "_extract_pypdf", return_value=("ab", 2)), \
         patch.object(skill, "_extract_pdfminer", return_value="full pdfminer text") as mock_pm:
        result = await skill.extract("dummy.pdf")
    mock_pm.assert_called_once()
    assert result.text == "full pdfminer text"


@pytest.mark.asyncio
async def test_pdf_skill_pdfminer_fallback_skipped_when_not_better(tmp_path):
    """PdfSkill does not replace pypdf text when pdfminer yields less."""
    from unittest.mock import patch
    from synthadoc.skills.pdf.scripts.main import PdfSkill

    skill = PdfSkill()
    good_text = "x" * 200
    with patch.object(skill, "_extract_pypdf", return_value=(good_text, 2)), \
         patch.object(skill, "_extract_pdfminer", return_value="short") as mock_pm:
        result = await skill.extract("dummy.pdf")
    mock_pm.assert_not_called()
    assert result.text == good_text


@pytest.mark.asyncio
async def test_pdf_skill_pdfminer_exception_returns_empty():
    """_extract_pdfminer returns empty string on any exception."""
    from unittest.mock import patch
    from synthadoc.skills.pdf.scripts.main import PdfSkill

    skill = PdfSkill()
    with patch("pdfminer.high_level.extract_text", side_effect=RuntimeError("boom")):
        result = skill._extract_pdfminer("dummy.pdf")
    assert result == ""


@pytest.mark.asyncio
async def test_xlsx_skill_extracts_sheets(tmp_path):
    """XlsxSkill reads sheet names and row data from a real xlsx file."""
    import openpyxl
    from synthadoc.skills.xlsx.scripts.main import XlsxSkill

    path = tmp_path / "data.xlsx"
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Report"
    ws.append(["Name", "Score"])
    ws.append(["Alice", 95])
    ws.append(["Bob", 87])
    wb.save(str(path))

    result = await XlsxSkill().extract(str(path))
    assert "Report" in result.text
    assert "Alice" in result.text
    assert "Bob" in result.text
    assert result.metadata["sheets"] == 1


@pytest.mark.asyncio
async def test_docx_skill_extracts_paragraphs(tmp_path):
    """DocxSkill returns paragraph text from a real docx file."""
    from docx import Document
    from synthadoc.skills.docx.scripts.main import DocxSkill

    path = tmp_path / "report.docx"
    doc = Document()
    doc.add_paragraph("Introduction paragraph.")
    doc.add_paragraph("Second paragraph with details.")
    doc.save(str(path))

    result = await DocxSkill().extract(str(path))
    assert "Introduction paragraph." in result.text
    assert "Second paragraph" in result.text
    assert result.metadata["paragraphs"] >= 2
